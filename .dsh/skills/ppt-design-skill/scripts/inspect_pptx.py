"""Print a compact structural report for a PPTX artifact."""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation


if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")


def inspect(path: Path) -> dict:
    prs = Presentation(str(path))
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = []
        kinds = {}
        for shape in slide.shapes:
            kind = str(shape.shape_type)
            kinds[kind] = kinds.get(kind, 0) + 1
            if getattr(shape, "has_text_frame", False):
                value = shape.text.strip()
                if value:
                    texts.append(value[:160])
        slides.append({"slide": index, "shapes": len(slide.shapes), "shape_types": kinds, "texts": texts})
    return {
        "path": str(path),
        "slide_count": len(prs.slides),
        "width_inches": round(prs.slide_width / 914400, 3),
        "height_inches": round(prs.slide_height / 914400, 3),
        "slides": slides,
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--pretty", action="store_true")
    args = parser.parse_args()
    report = inspect(args.pptx)
    print(json.dumps(report, ensure_ascii=False, indent=2 if args.pretty else None))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
